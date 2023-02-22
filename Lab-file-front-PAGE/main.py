import os
from pptx import Presentation
import openpyxl

def main():
    # Print menu options
    print("Menu:")
    print("1. Update PowerPoint presentation with data from Excel sheet")
    print("2. Exit")

    # Get user input for menu choice
    choice = input("Enter your choice (1 or 2): ")

    if choice == '1':
        # Set the path to the PowerPoint presentation file
        pptx_file_path = input("Enter the path to the PowerPoint presentation file: ")

        # Set the path to the Excel sheet
        excel_file_path = input("Enter the path to the Excel sheet: ")

        # Get the sheet name from the user
        sheet_name = input("Enter the name of the sheet containing the data: ")
        

        # Open the Excel sheet
        wb = openpyxl.load_workbook(excel_file_path)
        sheet = wb[sheet_name]

        # Loop through each rowumn in the Excel sheet
        for row in sheet.iter_rows(min_col=1):
            # Set the replacement text for the fields you want to replace
            replace_dict = {row[0].value: val.value for val in row[1:]}
            
            replace_dict = {
                'Name': row[3].value,
                'Subject': row[0].value,
                'Teacherame': row[1].value,
                'Teacher Position': row[2].value,
                'Roll No': row[4].value,
                'Section': row[5].value
            }

            # Open the PowerPoint presentation
            prs = Presentation(pptx_file_path)

            # Loop through each slide in the presentation
            for slide in prs.slides:
                # Loop through each shape in the slide
                for shape in slide.shapes:
                    # Check if the shape is a text box
                    if shape.has_text_frame:
                        # Loop through each paragraph in the text box
                        for paragraph in shape.text_frame.paragraphs:
                            # Loop through each run in the paragraph
                            for run in paragraph.runs:
                                # Loop through each key-value pair in the replace_dict
                                for key, value in replace_dict.items():
                                    # Replace the key with the value in the run text
                                    run.text = run.text.replace(str(key), str(value))

            # Save the modified PowerPoint presentation
            modified_pptx_file_path = f"{row[0].value}.pptx"
            prs.save(modified_pptx_file_path)

        # Close the Excel sheet
        wb.close()

        # Print success message
        print("PowerPoint presentations have been updated.")

    elif choice == '2':
        # Exit the program
        print("Exiting program...")
        return

    else:
        # Invalid input
        print("Invalid choice. Please enter 1 or 2.")
        main()

if __name__ == '__main__':
    main()
